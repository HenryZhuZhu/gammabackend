import os
import io
from typing import List, Dict, Any, Optional

import requests
from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pptx import Presentation


# ================= 配置区域 =================

GAMMA_API_KEY = os.getenv("GAMMA_API_KEY")
GAMMA_TEMPLATE_ID = os.getenv("GAMMA_TEMPLATE_ID")
GAMMA_THEME_ID = os.getenv("GAMMA_THEME_ID")
GAMMA_FOLDER_IDS = os.getenv("GAMMA_FOLDER_IDS")
GAMMA_EXPORT_FORMAT = os.getenv("GAMMA_EXPORT_AS", "pdf")  # "pdf" 或 "pptx"

GAMMA_BASE_URL = "https://public-api.gamma.app/v1.0"

app = FastAPI(title="AIStoryteller Gamma Backend (Template Mode)")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 可以后面改成你的 Netlify 域名
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# =============== 工具函数：从 PPT 提取文本 ===============

def extract_ppt_structure_and_text(file_bytes: bytes) -> Dict[str, Any]:
    """
    输入 PPTX 二进制，输出：
      - slides: 用于前端预览的结构
      - outline_text: 传给 Gamma 的文本大纲
    """
    try:
        from io import BytesIO
        prs = Presentation(BytesIO(file_bytes))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to open PPTX: {e}")

    slides_data: List[Dict[str, Any]] = []
    outline_parts: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        shapes_data = []
        slide_text_lines = []

        for shape in slide.shapes:
            text = ""
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                paragraphs = [p.text for p in shape.text_frame.paragraphs if p.text]
                text = "\n".join(paragraphs).strip()

            if text:
                shapes_data.append({"text": text})
                slide_text_lines.append(text)

        slides_data.append(
            {
                "index": idx,
                "shapes": shapes_data,
            }
        )

        if slide_text_lines:
            outline_parts.append(f"Slide {idx}:\n" + "\n".join(slide_text_lines))
        else:
            outline_parts.append(f"Slide {idx}:\n(No visible text content)")

    outline_text = "\n---\n".join(outline_parts)

    return {
        "slides": slides_data,
        "outline_text": outline_text,
    }


# =============== Gamma API 封装 ===============

def call_gamma_from_template(prompt_text: str) -> str:
    """
    调用 Gamma Create-from-template：
      POST /v1.0/generations/from-template
    返回 generationId
    """
    if not GAMMA_API_KEY:
        raise HTTPException(
            status_code=500,
            detail="GAMMA_API_KEY is not set in environment variables.",
        )

    if not GAMMA_TEMPLATE_ID:
        raise HTTPException(
            status_code=500,
            detail="GAMMA_TEMPLATE_ID is not set in environment variables.",
        )

    url = f"{GAMMA_BASE_URL}/generations/from-template"

    folder_ids = None
    if GAMMA_FOLDER_IDS:
        folder_ids = [f.strip() for f in GAMMA_FOLDER_IDS.split(",") if f.strip()]

    payload: Dict[str, Any] = {
        "gammaId": GAMMA_TEMPLATE_ID,
        "prompt": prompt_text,
        "exportAs": GAMMA_EXPORT_FORMAT,  # "pdf" or "pptx"
    }

    if GAMMA_THEME_ID:
        payload["themeId"] = GAMMA_THEME_ID

    if folder_ids:
        payload["folderIds"] = folder_ids

    headers = {
        "Content-Type": "application/json",
        "X-API-KEY": GAMMA_API_KEY,
    }

    try:
        resp = requests.post(url, json=payload, headers=headers, timeout=60)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to call Gamma: {e}")

    # ✅ 任何 2xx 都算成功
    if not resp.ok:
        raise HTTPException(
            status_code=resp.status_code,
            detail=f"Gamma create-from-template error: {resp.text}",
        )

    try:
        data = resp.json()
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Gamma response is not valid JSON: {e}, raw: {resp.text}",
        )

    generation_id = data.get("generationId")
    if not generation_id:
        raise HTTPException(
            status_code=500,
            detail=f"Gamma did not return generationId. Raw response: {data}",
        )

    return generation_id


def get_gamma_generation(generation_id: str) -> Dict[str, Any]:
    """
    GET /v1.0/generations/{generationId}
    返回 Gamma 的生成状态 + fileUrls + gammaUrl 等
    """
    if not GAMMA_API_KEY:
        raise HTTPException(
            status_code=500,
            detail="GAMMA_API_KEY is not set in environment variables.",
        )

    url = f"{GAMMA_BASE_URL}/generations/{generation_id}"
    headers = {
        "X-API-KEY": GAMMA_API_KEY,
        "accept": "application/json",
    }

    try:
        resp = requests.get(url, headers=headers, timeout=30)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to poll Gamma: {e}")

    if not resp.ok:
        raise HTTPException(
            status_code=resp.status_code,
            detail=f"Gamma status error: {resp.text}",
        )

    try:
        data = resp.json()
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Gamma status is not valid JSON: {e}, raw: {resp.text}",
        )

    return data


def download_gamma_file(gamma_result: Dict[str, Any]) -> bytes:
    """
    从 Gamma 结果中找到导出文件 URL 并下载。

    按 Gamma 最新返回结果格式：
      - 如果设置了 exportAs（pdf 或 pptx），响应里会有一个单独的 exportUrl 字段，
        比如：
        {
          "status": "completed",
          "gammaUrl": "...",
          "exportUrl": "https://assets.api.gamma.app/export/pdf/.../xxx.pdf",
          ...
        }
    所以我们优先使用 exportUrl。
    """
    file_url: Optional[str] = None

    # 1. 优先使用 exportUrl（这是文档和你日志里都实锤存在的字段）
    file_url = gamma_result.get("exportUrl")

    # 2. 向下兼容旧结构（如果将来 Gamma 再改回 fileUrls 也不至于挂）
    if not file_url:
        file_urls = gamma_result.get("fileUrls") or gamma_result.get("files") or {}
        if isinstance(file_urls, dict):
            # 这里不再用 expected_format，直接随缘取一个
            file_url = file_urls.get("pdf") or file_urls.get("pptx")

    # 3. 再兜底 pdfUrl/pptxUrl
    if not file_url:
        file_url = gamma_result.get("pdfUrl") or gamma_result.get("pptxUrl")

    if not file_url:
        raise HTTPException(
            status_code=500,
            detail=f"Gamma result completed but did not include exportUrl/file URL. Raw result: {gamma_result}",
        )

    try:
        resp = requests.get(file_url, timeout=120)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to download file from Gamma: {e}")

    if not resp.ok:
        raise HTTPException(
            status_code=resp.status_code,
            detail=f"Failed to download file from Gamma: {resp.text}",
        )

    return resp.content


# =============== API 路由 ===============

@app.get("/health")
def health_check():
    return {"status": "ok"}


@app.post("/api/parse_ppt")
async def parse_ppt(file: UploadFile = File(...)):
    """
    接收用户上传 PPTX，返回 slides 结构给前端预览
    """
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    file_bytes = await file.read()
    result = extract_ppt_structure_and_text(file_bytes)

    return JSONResponse({"slides": result["slides"]})


@app.post("/api/beautify_start")
async def beautify_start(file: UploadFile = File(...)):
    """
    步骤1：接收 PPTX，提取文本，调用 Gamma 创建任务，返回 generationId
    """
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    file_bytes = await file.read()
    parsed = extract_ppt_structure_and_text(file_bytes)
    outline_text = parsed["outline_text"]
  # 🔍 调试：打印从 PPT 中抽取出来的文字
    print("===== OUTLINE TEXT BEGIN =====")
    print(outline_text)
    print("=====  OUTLINE TEXT END  =====")
    prompt = (
    "You are an expert presentation designer. You are NOT creating new content, "
    "you are ONLY re-laying out an existing slide deck 1:1 using a Gamma template.\n\n"
    "SOURCE CONTENT:\n"
    "- The user PPTX has already been fully extracted into plain text below as 'Slide 1', 'Slide 2', etc.\n"
    "- This text is the ONLY source of truth. You MUST treat it as fixed, authoritative content.\n\n"
    "HARD, NON-NEGOTIABLE RULES (CONTENT):\n"
    "1. DO NOT change the topic or domain of the presentation under ANY circumstance.\n"
    "2. DO NOT add any new ideas, examples, numbers, claims, or narrative that are not present in the source.\n"
    "3. DO NOT remove any meaningful sentence, bullet, metric, or statement unless it is a literal duplicate.\n"
    "4. DO NOT paraphrase or rewrite technical/financial concepts into different wording. Keep the same meaning and terms.\n"
    "5. DO NOT infer or 'fill in' missing context. If something is not in the source, you must not invent it.\n"
    "6. DO NOT change company names, product names, ticker symbols, metrics, or numerical values.\n"
    "7. DO NOT change the language: Chinese text must stay in Chinese, English text must stay in English.\n"
    "8. DO NOT turn a stock analysis deck into a semiconductor, AI, or other domain deck unless the source itself is about that domain.\n"
    "9. USE the main title text from Slide 1 as the presentation title. Do not replace it with another theme.\n\n"
    "ALLOWED ACTIONS (VISUAL / LAYOUT ONLY):\n"
    "- Improve visual hierarchy (titles, subtitles, bullet levels, spacing, grouping).\n"
    "- Re-arrange blocks of text on each slide to make them clearer and easier to read.\n"
    "- Convert raw text into bullets or sections as long as the wording and meaning do not change.\n"
    "- Split a long slide into two slides ONLY if the content becomes more readable, but all text must still appear.\n"
    "- Merge two slides ONLY if they are obviously the same topic and ALL text is preserved.\n\n"
    "SLIDE-BY-SLIDE MAPPING RULES:\n"
    "- For each 'Slide N' block in the source, you MUST create a corresponding slide in the final deck.\n"
    "- The content appearing under 'Slide N:' in the source must all appear somewhere on the slide(s) corresponding to Slide N.\n"
    "- Do not move content from Slide N into a completely different topic slide.\n"
    "- The logical grouping per slide should stay as close as possible to the original deck.\n\n"
    "SELF-CHECK BEFORE FINALIZING:\n"
    "Before you finalize the deck, mentally verify these conditions:\n"
    "- Every sentence, bullet, or metric from the source text appears somewhere in the resulting slides.\n"
    "- No slide title has been changed to a different domain (for example, do NOT change a stock analysis title "
    "into a 'Semiconductor Testing & AI Product Analysis' title unless that is explicitly in the source).\n"
    "- No new sections or talking points were introduced that are not expressed in the source text.\n\n"
    "YOUR GOAL:\n"
    "- Produce a version of the same deck that looks cleaner, more professional, and better structured, "
    "but whose textual content could be line-by-line matched back to the original.\n\n"
    "Below is the extracted content from the user's PPTX. For each 'Slide X', treat its lines as the content "
    "you must faithfully preserve for that slide, only improving layout and visual presentation:\n\n"
    f"{outline_text}"
)


    generation_id = call_gamma_from_template(prompt)

    return JSONResponse({"generationId": generation_id})





@app.get("/api/beautify_status")
def beautify_status(generationId: str = Query(..., alias="generationId")):
    """
    步骤2：前端轮询调用，查询 Gamma 任务状态
    返回 { status, gammaUrl }
    """
    data = get_gamma_generation(generationId)
    status = data.get("status", "unknown")
    gamma_url = data.get("gammaUrl")

    return JSONResponse({"status": status, "gammaUrl": gamma_url})


@app.get("/api/beautify_result")
def beautify_result(
    generationId: str = Query(..., alias="generationId"),
    filename: Optional[str] = Query(None),
):
    """
    步骤3：任务完成后，下载最终 PDF/PPTX 并返回给前端
    """
    data = get_gamma_generation(generationId)
    status = data.get("status")

    if status != "completed":
        raise HTTPException(
            status_code=400,
            detail=f"Gamma generation is not completed yet. Current status: {status}",
        )

    file_bytes = download_gamma_file(data)

    base_name = os.path.splitext(filename or "presentation")[0]
    ext = "pdf" if GAMMA_EXPORT_FORMAT.lower() == "pdf" else "pptx"
    output_filename = f"{base_name}_gamma_beautified.{ext}"

    media_type = (
        "application/pdf"
        if ext == "pdf"
        else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    return StreamingResponse(
        io.BytesIO(file_bytes),
        media_type=media_type,
        headers={
            "Content-Disposition": f'attachment; filename="{output_filename}"'
        },
    )


# （可选）保留旧 /api/beautify，提示前端不要再用
@app.post("/api/beautify")
async def beautify_legacy():
    raise HTTPException(
        status_code=410,
        detail="Deprecated endpoint. Please use /api/beautify_start + /api/beautify_status + /api/beautify_result.",
    )
